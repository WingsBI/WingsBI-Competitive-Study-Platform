import os
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from skimage.metrics import structural_similarity as ssim
import numpy as np
from datetime import datetime

def compare_images(img1_path, img2_path):
    img1 = Image.open(img1_path).convert('L')  # Convert to grayscale
    img2 = Image.open(img2_path).convert('L')  # Convert to grayscale

    img1_array = np.array(img1)
    img2_array = np.array(img2)

    # Compute the SSIM between the two images
    score, _ = ssim(img1_array, img2_array, full=True)
    
    # Calculate percentage similarity
    percentage_similarity = score * 100

    # Return percentage similarity as a string
    return f"Similarity: {percentage_similarity:.2f}%"

def create_presentation(csv_file_path, wings_folder, zoho_folder, questions_df):
    # Initialize presentation object
    prs = Presentation()

    def add_title_slide(prs):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Competitive Study"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(72)
        title.text_frame.paragraphs[0].font.underline = True

        subtitle = slide.placeholders[1]
        subtitle.text = questions_df.iloc[2, 10]
        subtitle.text_frame.paragraphs[0].font.bold = True
        subtitle.text_frame.paragraphs[0].font.size = Pt(36)
        subtitle.text_frame.paragraphs[0].font.underline = True

                # Add timestamp
        # timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        # timestamp_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.5))
        # timestamp_frame = timestamp_box.text_frame
        # timestamp_frame.text = timestamp
        # timestamp_frame.paragraphs[0].font.size = Pt(12)
        # timestamp_frame.paragraphs[0].font.italic = True
        # timestamp_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    def add_image_slide(prs, image_path1, image_path2, title_text, slide_number):
        slide_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(slide_layout)
        # Add empty text box at the bottom of the slide
        left = Inches(0.5)  # Adjust left position as needed
        top = Inches(6.9)   # Adjust top position as needed
        width = Inches(9)   # Adjust width as needed
        height = Inches(0.4) # Adjust height as needed
        textbox = slide.shapes.add_textbox(left, top, width, height)

        # Compare the images and get result
        comparison_result = compare_images(image_path1, image_path2)

        # Format the text box with black border
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill color
        textbox.line.color.rgb = RGBColor(0, 0, 0)  # Black border color
        textbox.line.width = Pt(1)  # Border width

        # Replace underscores with question marks in title_text
        title_text = title_text.replace('_', '?')

        # Construct the title with prefix "Q.<slide_number>"
        slide_title = f"Q.{slide_number}. {title_text}"

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Add 'Wings' to the first placeholder (subtitle placeholder)
        subtitle_shape = slide.placeholders[1].text_frame
        subtitle_wings = subtitle_shape.add_paragraph()
        subtitle_wings.text = "Wings"
        subtitle_wings.font.size = Pt(20)
        subtitle_wings.font.bold = True
        subtitle_wings.alignment = PP_ALIGN.CENTER

        # Add image from wings folder
        if os.path.exists(image_path1):
            left = slide.placeholders[2].left
            top = slide.placeholders[2].top
            width = slide.placeholders[2].width
            height = slide.placeholders[2].height
            slide.shapes.add_picture(image_path1, left, top, width, height)

        # Add 'Zoho' to the second placeholder (adjacent placeholder)
        if len(slide.placeholders) > 3:  # Ensure there is a second placeholder
            zoho_shape = slide.placeholders[3].text_frame
            subtitle_zoho = zoho_shape.add_paragraph()
            subtitle_zoho.text = "Zoho"
            subtitle_zoho.font.size = Pt(20)
            subtitle_zoho.font.bold = True
            subtitle_zoho.alignment = PP_ALIGN.CENTER

        # Add image from zoho folder
        if os.path.exists(image_path2):
            left = slide.placeholders[4].left
            top = slide.placeholders[4].top
            width = slide.placeholders[4].width
            height = slide.placeholders[4].height
            slide.shapes.add_picture(image_path2, left, top, width, height)

        # Add comparison result to the text box
        textbox.text = comparison_result

    def add_images_from_folder(folder_path1, folder_path2, folder_name):
        files1 = os.listdir(folder_path1)
        files2 = os.listdir(folder_path2)

        # Assuming both folders have the same number of images
        num_images = min(len(files1), len(files2))

        for i in range(num_images):
            if files1[i].endswith(".png") and files2[i].endswith(".png"):
                image_path1 = os.path.join(folder_path1, files1[i])
                image_path2 = os.path.join(folder_path2, files2[i])
                title_text = os.path.splitext(files1[i])[0]  # Get the filename without extension
                slide_number = i + 1  # Slide numbers start from 1
                add_image_slide(prs, image_path1, image_path2, title_text, slide_number)

    add_title_slide(prs)
    add_images_from_folder(wings_folder, zoho_folder, "WingsBI and Zoho Analytics")

    # Determine the desktop path for the current user
    common_directory = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ppt_path = os.path.join(common_directory, "PPT")

    # Add timestamp to the filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_pptx_file = os.path.join(ppt_path, f'Competitive Study (Insurance Database)_{timestamp}.pptx')
    prs.save(output_pptx_file)

    print(f"Presentation saved as '{output_pptx_file}'")

    # Launch the file (Windows specific)
    os.startfile(output_pptx_file)
